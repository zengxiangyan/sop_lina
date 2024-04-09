# -*- coding: utf-8 -*-
from openpyxl import load_workbook
from datetime import datetime
from dateutil.relativedelta import relativedelta
from pptx.dml.color import ColorFormat, RGBColor
import re

import myoffice
from pptx import Presentation
import time
from clickhouse_sqlalchemy import make_session
from sqlalchemy import create_engine
import pandas as pd
import concurrent.futures

def config_date(report_date):
    date1 = str(datetime.strptime(report_date,"%Y-%m-%d").year-2) + '-' + str(datetime.strptime(report_date,"%Y-%m-%d").month) + '-' + str(datetime.strptime(report_date,"%Y-%m-%d").day)
    date2 = str(datetime.strptime(report_date,"%Y-%m-%d").year) + '-' + str(datetime.strptime(report_date,"%Y-%m-%d").month) + '-' + str(datetime.strptime(report_date,"%Y-%m-%d").day)
    date = f"date >= '{date1}' AND date < '{date2}'"
    sql_list = [f"""SELECT toStartOfMonth(pkey) AS Gmonth,
    CASE WHEN 1 THEN '行业|销售额(000元)' END AS `大盘类型`,
    CASE WHEN `source` = 1 AND shop_type in (23,26) THEN '国内猫超' 
         WHEN `source` = 1 AND shop_type in (21,25) THEN '国内天猫(旗舰店+分销店)'
         WHEN `source` = 2 AND shop_type in (11,12) THEN '京东(国内自营+国内POP)'
         ELSE '其它'
         END AS `渠道`,`sp厂商`,`sp子品类`,`sp产品品牌`,`sp子品牌`,`sp件数`,item_id,trade_props.name,trade_props.value as `交易属性`,
    sum(if(model!=0,sales/100,0))/(model*sum(if(model!=0,num,0)))*1000 as p4,
    toFloat64OrZero(`sp总规格`) as model,
    `sp是否人工`,`sp适用人群（段数）`,SUM(num) AS total_num, SUM(sales)/100 AS total_sales  
    FROM sop_e.entity_prod_91357_E_meisu WHERE ({date} 
    AND source*100+shop_type IN [123,126,121,125,211,212] 
    AND cid in (211104,201284105,7052,31762) 
    AND `sp是否无效链接`!='无效链接') GROUP BY `渠道`,Gmonth,item_id,`sp件数`,model,`sp子品类`,`sp是否人工`,`sp适用人群（段数）`,`sp厂商`,`sp产品品牌`,`sp子品牌`,`交易属性`,trade_props.name
    """,
    f"""SELECT toStartOfMonth(pkey) AS Gmonth,
    CASE WHEN 1 THEN '美素佳儿（金装+源悦+皇家）' END AS `大盘类型`,
    CASE WHEN `source` = 1 AND shop_type in (23,26) THEN '国内猫超' 
         WHEN `source` = 1 AND shop_type in (21,25) THEN '国内天猫(旗舰店+分销店)'
         WHEN `source` = 2 AND shop_type in (11,12) THEN '京东(国内自营+国内POP)'
         ELSE '其它'
         END AS `渠道`,`sp厂商`,`sp子品类`,`sp产品品牌`,`sp子品牌`,`sp件数`,item_id,trade_props.name,trade_props.value as `交易属性`,
    sum(if(model!=0,sales/100,0))/(model*sum(if(model!=0,num,0)))*1000 as p4,
    toFloat64OrZero(`sp总规格`) as model,
    `sp是否人工`,`sp适用人群（段数）`,SUM(num) AS total_num, SUM(sales)/100 AS total_sales 
    FROM sop_e.entity_prod_91357_E_meisu WHERE ({date} 
    AND source*100+shop_type IN [123,126,121,125,211,212] 
    AND cid in (211104,201284105,7052,31762) 
    AND `sp厂商`='Frieslandcampina/荷兰皇家菲仕兰' 
    AND `sp是否无效链接`!='无效链接') GROUP BY `渠道`,Gmonth,item_id,`sp件数`,model,`sp子品类`,`sp是否人工`,`sp适用人群（段数）`,`sp厂商`,`sp产品品牌`,`sp子品牌`,`交易属性`,trade_props.name
    """,
    f"""SELECT toStartOfMonth(pkey) AS Gmonth,
    CASE WHEN 1 THEN '皇家美素佳儿' END AS `大盘类型`,
    CASE WHEN `source` = 1 AND shop_type in (23,26) THEN '国内猫超' 
         WHEN `source` = 1 AND shop_type in (21,25) THEN '国内天猫(旗舰店+分销店)'
         WHEN `source` = 2 AND shop_type in (11,12) THEN '京东(国内自营+国内POP)'
         ELSE '其它'
         END AS `渠道`,`sp厂商`,`sp子品类`,`sp产品品牌`,`sp子品牌`,`sp件数`,item_id,trade_props.name,trade_props.value as `交易属性`,
    sum(if(model!=0,sales/100,0))/(model*sum(if(model!=0,num,0)))*1000 as p4,
    toFloat64OrZero(`sp总规格`) as model,
    `sp是否人工`,`sp适用人群（段数）`,SUM(num) AS total_num, SUM(sales)/100 AS total_sales 
    FROM sop_e.entity_prod_91357_E_meisu WHERE ({date} 
    AND source*100+shop_type IN [123,126,121,125,211,212] 
    AND cid in (211104,201284105,7052,31762) 
    AND `sp产品品牌`='Friso Prestige/皇家美素佳儿'
    AND `sp子品牌`!='旺玥'
    AND `sp是否无效链接`!='无效链接') GROUP BY `渠道`,Gmonth,item_id,`sp件数`,model,`sp子品类`,`sp是否人工`,`sp适用人群（段数）`,`sp厂商`,`sp产品品牌`,`sp子品牌`,`交易属性`,trade_props.name
    """,
    f"""SELECT toStartOfMonth(pkey) AS Gmonth,
    CASE WHEN 1 THEN '美素佳儿（金装）' END AS `大盘类型`,
    CASE WHEN `source` = 1 AND shop_type in (23,26) THEN '国内猫超' 
         WHEN `source` = 1 AND shop_type in (21,25) THEN '国内天猫(旗舰店+分销店)'
         WHEN `source` = 2 AND shop_type in (11,12) THEN '京东(国内自营+国内POP)'
         ELSE '其它'
         END AS `渠道`,`sp厂商`,`sp子品类`,`sp产品品牌`,`sp子品牌`,`sp件数`,item_id,trade_props.name,trade_props.value as `交易属性`,
    sum(if(model!=0,sales/100,0))/(model*sum(if(model!=0,num,0)))*1000 as p4,
    toFloat64OrZero(`sp总规格`) as model,
    `sp是否人工`,`sp适用人群（段数）`,SUM(num) AS total_num, SUM(sales)/100 AS total_sales 
    FROM sop_e.entity_prod_91357_E_meisu WHERE ({date} 
    AND source*100+shop_type IN [123,126,121,125,211,212] 
    AND cid in (211104,201284105,7052,31762)  
    AND `sp产品品牌`='Friso/美素佳儿'
    AND `sp子品牌`!='源悦'
    AND `sp是否无效链接`!='无效链接') GROUP BY `渠道`,Gmonth,item_id,`sp件数`,model,`sp子品类`,`sp是否人工`,`sp适用人群（段数）`,`sp厂商`,`sp产品品牌`,`sp子品牌`,`交易属性`,trade_props.name
    """,
    f"""SELECT toStartOfMonth(pkey) AS Gmonth,
    CASE WHEN 1 THEN '美素佳儿（源悦）' END AS `大盘类型`,
    CASE WHEN `source` = 1 AND shop_type in (23,26) THEN '国内猫超' 
         WHEN `source` = 1 AND shop_type in (21,25) THEN '国内天猫(旗舰店+分销店)'
         WHEN `source` = 2 AND shop_type in (11,12) THEN '京东(国内自营+国内POP)'
         ELSE '其它'
         END AS `渠道`,`sp厂商`,`sp子品类`,`sp产品品牌`,`sp子品牌`,`sp件数`,item_id,trade_props.name,trade_props.value as `交易属性`,
    sum(if(model!=0,sales/100,0))/(model*sum(if(model!=0,num,0)))*1000 as p4,
    toFloat64OrZero(`sp总规格`) as model,
    `sp是否人工`,`sp适用人群（段数）`,SUM(num) AS total_num, SUM(sales)/100 AS total_sales 
    FROM sop_e.entity_prod_91357_E_meisu WHERE ({date} 
    AND source*100+shop_type IN [123,126,121,125,211,212] 
    AND cid in (211104,201284105,7052,31762)  
    AND `sp子品牌`='源悦'
    AND `sp是否无效链接`!='无效链接') GROUP BY `渠道`,Gmonth,item_id,`sp件数`,model,`sp子品类`,`sp是否人工`,`sp适用人群（段数）`,`sp厂商`,`sp产品品牌`,`sp子品牌`,`交易属性`,trade_props.name
    """,
    f"""SELECT toStartOfMonth(pkey) AS Gmonth,
    CASE WHEN 1 THEN '旺玥' END AS `大盘类型`,
    CASE WHEN `source` = 1 AND shop_type in (23,26) THEN '国内猫超' 
         WHEN `source` = 1 AND shop_type in (21,25) THEN '国内天猫(旗舰店+分销店)'
         WHEN `source` = 2 AND shop_type in (11,12) THEN '京东(国内自营+国内POP)'
         ELSE '其它'
         END AS `渠道`,`sp厂商`,`sp子品类`,`sp产品品牌`,`sp子品牌`,`sp件数`,item_id,trade_props.name,trade_props.value as `交易属性`,
    sum(if(model!=0,sales/100,0))/(model*sum(if(model!=0,num,0)))*1000 as p4,
    toFloat64OrZero(`sp总规格`) as model,
    `sp是否人工`,`sp适用人群（段数）`,SUM(num) AS total_num, SUM(sales)/100 AS total_sales 
    FROM sop_e.entity_prod_91357_E_meisu WHERE ({date} 
    AND source*100+shop_type IN [123,126,121,125,211,212] 
    AND cid in (211104,201284105,7052,31762) 
    AND `sp产品品牌`='Friso Prestige/皇家美素佳儿'
    AND `sp子品牌`='旺玥'
    AND `sp是否无效链接`!='无效链接') GROUP BY `渠道`,Gmonth,item_id,`sp件数`,model,`sp子品类`,`sp是否人工`,`sp适用人群（段数）`,`sp厂商`,`sp产品品牌`,`sp子品牌`,`交易属性`,trade_props.name
    """]
    return date1,date2,sql_list


def chsql_db_connect():
    conf = {
        "user": "yinglina",
        "password": "xfUW5GMr",
        "server_host": "127.0.0.1",
        "port": "10192",
        "db": "sop"
    }
    connection = 'clickhouse://{user}:{password}@{server_host}:{port}/{db}'.format(**conf)
    engine = create_engine(connection, pool_size=100, pool_recycle=3600, pool_timeout=20)
    return engine
def execute_sql(sql):
    engine = chsql_db_connect()
    session = make_session(engine)
    cursor = session.execute(sql)
    try:
        fields = cursor._metadata.keys
        mydata = pd.DataFrame([dict(zip(fields, item)) for item in cursor.fetchall()])
        return mydata
    finally:
        cursor.close()
        session.close()

def importdata_new(sql_list):
    # 使用线程池执行 SQL 查询
    with concurrent.futures.ThreadPoolExecutor() as executor:
        futures = [executor.submit(execute_sql, sql) for sql in sql_list]
        results = [future.result() for future in concurrent.futures.as_completed(futures)]
    # 将结果合并
    data = pd.concat(results, axis=0, ignore_index=True)
    data.loc[data['sp适用人群（段数）']=='5段','sp适用人群（段数）'] ='其它'
    return data

def importdata(sql_list):
    conf = {
        "user": "yinglina",
        "password": "xfUW5GMr",
        "server_host": "10.21.90.15",
        "port": "10081",
        "db": "sop"
    }
    connection = 'clickhouse://{user}:{password}@{server_host}:{port}/{db}'.format(**conf)
    engine = create_engine(connection, pool_size=100, pool_recycle=3600, pool_timeout=20)
    # sql = 'select * from sop_e.entity_prod_92087_E limit 1,1000'
    for num, sql in enumerate(sql_list):
        print(num)
        session = make_session(engine)
        cursor = session.execute(sql)
        try:
            fields = cursor._metadata.keys
            mydata = pd.DataFrame([dict(zip(fields, item)) for item in cursor.fetchall()])
        finally:
            cursor.close()
            session.close()
        if num == 0:
            newdata = mydata
        else:
            newdata = pd.concat([newdata, mydata], axis=0, ignore_index=True)
    return newdata


def get_source_model(report_date):
    month_dic = {'1': 'Jan', '2': 'Feb', '3': 'Mar', '4': 'Apr', '5': 'May', '6': 'Jun', '7': 'Jul', '8': 'Aug',
                 '9': 'Sep', '10': 'Oct', '11': 'Nov', '12': 'Dec'}
    tt = report_date
    table = load_workbook(project_path + table_template + '.xlsx')
    sheet = table['summary']
    for i in range(len(sql_list)):
        sheet.cell(1 + 7 * i, 2).value = (datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=3)).strftime('%Y%m')
        sheet.cell(1 + 7 * i, 9).value = (datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=3)).strftime('%Y%m')
        sheet.cell(1 + 7 * i, 3).value = (datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=2)).strftime('%Y%m')
        sheet.cell(1 + 7 * i, 10).value = (datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=2)).strftime('%Y%m')
        sheet.cell(1 + 7 * i, 4).value = (datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).strftime('%Y%m')
        sheet.cell(1 + 7 * i, 11).value = (datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).strftime('%Y%m')

    sheet = table['title']
    sheet.cell(1, 1).value = month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=13)).year)[-2:]
    sheet.cell(1, 2).value = 'MAT\n' + month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=13)).year)[-2:]
    sheet.cell(2, 1).value = month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).year)[-2:]
    sheet.cell(2, 2).value = 'MAT\n' + month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).year)[-2:]

    sheet = table['价格段1']
    sheet.cell(2, 1).value = 'MAT ' + month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=13)).year)[-2:]
    sheet.cell(3, 1).value = 'MAT ' + month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).year)[-2:]

    sheet = table['价格段2']
    sheet.cell(1, 2).value = 'MAT ' + month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=13)).year)[-2:]
    sheet.cell(1, 3).value = 'MAT ' + month_dic[str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).month)] + str((datetime.strptime(tt, '%Y-%m-%d') - relativedelta(months=1)).year)[-2:]
    table.save(project_path + table_template + '.xlsx')


def process_data():
    month = [m.strftime("%Y%m") for m in newdata['Gmonth']]
    newdata.loc[newdata['sp适用人群（段数）'] == '', 'sp适用人群（段数）'] = '其它'
    mat = []
    Sale_volume = []
    Avg_Price_kg = []
    for m in month:
        if int(m) >= 100 * datetime.strptime(date1, "%Y-%m-%d").year + datetime.strptime(date1,
                                                                                         "%Y-%m-%d").month and int(
                m) < 100 * (datetime.strptime(date2, "%Y-%m-%d").year - 1) + datetime.strptime(date2, "%Y-%m-%d").month:
            mat.append(mat_list[0])
        else:
            mat.append(mat_list[1])
    for i, num in enumerate(newdata['total_num']):
        try:
            value = float(newdata['model'][i])
            Sale_volume.append(value * num)
        except:
            print(newdata['model'][i], num)
            Sale_volume.append(0)
    for i, num in enumerate(newdata['p4']):
        if newdata['p4'][i] != '':
            if num < 210:
                Avg_Price_kg.append('M<210')
            elif (num >= 210 and num < 290):
                Avg_Price_kg.append('P[210,290)')
            elif (num >= 290 and num < 390):
                Avg_Price_kg.append('SP[290,390)')
            elif num >= 390:
                Avg_Price_kg.append('UP>=390')
            else:
                Avg_Price_kg.append('其它')
                # print("总规格", newdata['model'][i], "销售额", newdata['total_sales'][i], "销量", num) #剩余的都是总规格为0或者销量为0的链接无法区分价格段
        else:
            Avg_Price_kg.append('其它')
    newdata['Gmonth'] = month
    newdata['MAT类别'] = mat
    newdata['Sale volume'] = Sale_volume
    newdata['Avg_Price_kg'] = Avg_Price_kg
    return newdata


def get_summary():
    source = ['京东+天猫', '京东(国内自营+国内POP)', '国内天猫整体', '国内天猫(旗舰店+分销店)', '国内猫超']
    sheet = table['summary']
    for i, ms in enumerate(meisu):
        for j, m in enumerate(mouth):
            for k, s in enumerate(source):
                if s == '京东+天猫':
                    data = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == ms)]
                    data_total = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == meisu[0])]
                    last_year = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == ms)]
                    last_year_total = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == meisu[0])]
                if s == '国内天猫整体':
                    data = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == ms) & (newdata['渠道'] != '京东(国内自营+国内POP)')]
                    data_total = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == meisu[0]) & (newdata['渠道'] != '京东(国内自营+国内POP)')]
                    last_year = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == ms) & (newdata['渠道'] != '京东(国内自营+国内POP)')]
                    last_year_total = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == meisu[0]) & (newdata['渠道'] != '京东(国内自营+国内POP)')]
                if s in ['京东(国内自营+国内POP)', '国内天猫(旗舰店+分销店)', '国内猫超']:
                    data = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == ms) & (newdata['渠道'] == s)]
                    data_total = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == meisu[0]) & (newdata['渠道'] == s)]
                    last_year = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == ms) & (newdata['渠道'] == s)]
                    last_year_total = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == meisu[0]) & (newdata['渠道'] == s)]
                if ms == meisu[0]:
                    data_total = newdata.loc[(newdata['Gmonth'] == m) & (newdata['大盘类型'] == meisu[0])]
                    last_year_total = newdata.loc[(newdata['Gmonth'] == str(int(m) - 100)) & (newdata['大盘类型'] == meisu[0])]
                sheet.cell(2 + k + 7 * i, 2 + j).value = int("%.0f" % (data['total_sales'].sum() / 1000))
                value = data['total_sales'].sum() / data_total['total_sales'].sum() * 100
                sheet.cell(2 + k + 7 * i, 2 + j + 7).value = str("%.1f" % value) + "%"
                if j == len(mouth) - 1:
                    if sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 2).value != 0:
                        MOMGR = (sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 1).value - sheet.cell(2 + k + 7 * i,2 + len(mouth) - 2).value) / sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 2).value * 100
                        if str("%.1f" % MOMGR) not in ['inf', 'nan']:
                            sheet.cell(2 + k + 7 * i, 2 + len(mouth)).value = str("%.1f" % MOMGR) + "%"
                        MOM = float(re.findall(r'(\d+.\d+)', sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 1 + 7).value)[0]) - float(re.findall(r'(\d+.\d+)', sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 2 + 7).value)[0])
                        if str("%.1f" % MOM) not in ['inf', 'nan']:
                            sheet.cell(2 + k + 7 * i, 2 + len(mouth) + 7).value = str("%.1f" % MOM) + "%"
                    if last_year['total_sales'].sum() != 0:
                        YOYGR = (sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 1).value * 1000 - last_year['total_sales'].sum()) / last_year['total_sales'].sum() * 100
                        if str("%.1f" % YOYGR) not in ['inf', 'nan']:
                            sheet.cell(2 + k + 7 * i, 2 + len(mouth) + 1).value = str("%.1f" % YOYGR) + "%"
                        YOY = float(re.findall(r'(\d+.\d+)', sheet.cell(2 + k + 7 * i, 2 + len(mouth) - 1 + 7).value)[0]) - (last_year['total_sales'].sum() / (last_year_total['total_sales'].sum()) * 100)
                        if str("%.1f" % YOY) not in ['inf', 'nan']:
                            sheet.cell(2 + k + 7 * i, 2 + len(mouth) + 1 + 7).value = str("%.1f" % YOY) + "%"
    return table


def get_table_data():
    meisu = ['行业|销售额(000元)', '美素佳儿（金装+源悦+皇家）', '皇家美素佳儿', '旺玥', '美素佳儿（金装）','美素佳儿（源悦）']
    sources = {'渠道': ['京东(国内自营+国内POP)', '国内天猫(旗舰店+分销店)', '国内猫超'],
               'Avg_Price_kg': ['其它', 'UP>=390', 'SP[290,390)', 'P[210,290)', 'M<210'],
               'sp适用人群（段数）': ['1段', '2段', '3段', '4段', '其它']}
    sheetnames_list = [{'value YoY GR%': 'total_sales', 'volume YoY GR%': 'Sale volume'},
                       {'price value YoY GR%': 'total_sales', 'price volume YoY GR%': 'Sale volume'},
                       {'peo value YoY GR%': 'total_sales', 'peo volume YoY GR%': 'Sale volume'}]
    for num, pkey in enumerate(sources):
        sheetnames = sheetnames_list[num]
        source = sources[pkey]
        print(pkey)
        for sheet_name in list(sheetnames.keys()):
            sheet = table[sheet_name]
            for i, ms in enumerate(meisu):
                for k, s in enumerate(source):
                    data = newdata.loc[(newdata['MAT类别'] == mat_list[1]) & (newdata['大盘类型'] == ms) & (newdata[pkey] == s)]
                    last_year = newdata.loc[(newdata['MAT类别'] == mat_list[0]) & (newdata['大盘类型'] == ms) & (newdata[pkey] == s)]
                    value = get_percent((data[sheetnames[sheet_name]].sum() - last_year[sheetnames[sheet_name]].sum()),last_year[sheetnames[sheet_name]].sum())
                    if value != '':
                        if value > 0:
                            sheet.cell(1 + k, 1 + i).value = '+' + str("%.1f" % value) + "%"
                        if value < 0:
                            sheet.cell(1 + k, 1 + i).value = str("%.1f" % value) + "%"
    return table

def get_percent(numerator,denominator):
    if pd.notna(denominator) and denominator != 0:
        value = numerator / denominator * 100
    else:
        value = ''
    return value
def get_chart_data():
    global price_list, group_people
    # price_list = [list(newdata.groupby(['Avg_Price_kg'], as_index=False)['total_num', 'total_sales'].sum()['Avg_Price_kg'])[len(list(newdata.groupby(['Avg_Price_kg'], as_index=False)['total_num', 'total_sales'].sum()['Avg_Price_kg'])) - i - 1] for i in range(len(list(newdata.groupby(['Avg_Price_kg'], as_index=False)['total_num', 'total_sales'].sum()['Avg_Price_kg'])))]
    grouped_data = newdata.groupby('Avg_Price_kg', as_index=False)[['total_num', 'total_sales']].sum()
    grouped_data_sorted = grouped_data.sort_values('Avg_Price_kg', ascending=False)
    price_list = list(grouped_data_sorted['Avg_Price_kg'])

    group_people = list(newdata.groupby(['sp适用人群（段数）'], as_index=False)[['total_num', 'total_sales']].sum()['sp适用人群（段数）'])
    series_dict = {'渠道': [['京东(国内自营+国内POP)', '国内天猫(旗舰店+分销店)', '国内猫超'],{'Sales channel value M': 'total_sales', 'Sales channel volume': 'Sale volume','t1': 'total_sales', 't2': 'Sale volume'}],
                   'Avg_Price_kg': [price_list, {'Sales price value M': 'total_sales', 'Sales price volume': 'Sale volume'}],
                   'sp适用人群（段数）': [group_people,{'Sales people_group value M': 'total_sales','Sales people_group volume': 'Sale volume'}]}
    for series_name in series_dict.keys():
        series = series_dict[series_name][0]
        sheetnames = series_dict[series_name][1]
        for sheet_name in list(sheetnames.keys()):
            sheet = table[sheet_name]
            for i, s in enumerate(series):
                num = 0
                for j, ms in enumerate(meisu):
                    for k, mt in enumerate(mat_list):
                        data = newdata.loc[(newdata['MAT类别'] == mt) & (newdata['大盘类型'] == ms) & (newdata[series_name] == s)]
                        total_data = newdata.loc[(newdata['MAT类别'] == mt) & (newdata['大盘类型'] == ms)]
                        if sheet_name in ['t1', 't2']:
                            per = get_percent(data[sheetnames[sheet_name]].sum() , total_data[sheetnames[sheet_name]].sum())
                            if per == '':
                                sheet.cell(2 + k + j * 2 + num, 2 + i).value = ''
                            else:
                                sheet.cell(2 + k + j * 2 + num, 2 + i).value = "%.1f" % (per)
                        else:
                            sheet.cell(2 + k + j * 2 + num, 2 + i).value = get_percent(data[sheetnames[sheet_name]].sum() , total_data[sheetnames[sheet_name]].sum())
                        if int(total_data[sheetnames[sheet_name]].sum() / 1000000000) != 0:
                            val0 = "%.0f" % (total_data[sheetnames[sheet_name]].sum() / 1000000)
                            val = change_int_type(val0)
                            sheet.cell(2 + k + j * 2 + num, 2 + len(series)).value = val
                        elif int(total_data[sheetnames[sheet_name]].sum() / 1000000) != 0:
                            sheet.cell(2 + k + j * 2 + num, 2 + len(series)).value = "%.0f" % (total_data[sheetnames[sheet_name]].sum() / 1000000)
                        else:
                            sheet.cell(2 + k + j * 2 + num, 2 + len(series)).value = "%.2f" % (total_data[sheetnames[sheet_name]].sum() / 1000000)
                    num += 1
    return table


def total_price_chart():
    sheet = table['价格段1']
    for i, s in enumerate(price_list):
        for j, mt in enumerate(mat_list):
            data = newdata.loc[(newdata['MAT类别'] == mt) & (newdata['大盘类型'] == '行业|销售额(000元)') & (newdata['Avg_Price_kg'] == s)]
            total_data = newdata.loc[(newdata['MAT类别'] == mt) & (newdata['大盘类型'] == '行业|销售额(000元)')]
            sheet.cell(2 + j, 2 + i).value = data['total_sales'].sum() / total_data['total_sales'].sum() * 100
            if int(total_data['total_sales'].sum() / 1000000000) != 0:
                val0 = "%.0f" % (total_data['total_sales'].sum() / 1000000)
                val = change_int_type(val0)
                sheet.cell(2 + j, 2 + len(price_list)).value = val
            if int(total_data['total_sales'].sum() / 1000000000) == 0:
                sheet.cell(2 + j, 2 + len(price_list)).value = "%.0f" % (total_data['total_sales'].sum() / 1000000)
    return table


def total_price_table():
    sheet = table['价格段2']
    for i, mt in enumerate(mat_list):
        for j, s in enumerate(price_list):
            data = newdata.loc[(newdata['MAT类别'] == mt) & (newdata['大盘类型'] == '行业|销售额(000元)') & (newdata['Avg_Price_kg'] == s)]
            # total_data = newdata.loc[(newdata['MAT类别'] == mt) & (newdata['大盘类型'] == '行业|销售额(000元)') & (newdata['sp是否人工'] != '否')] #这张表不用排除非人工答题
            per = get_percent(data['total_sales'].sum() * 10 , data['Sale volume'].sum())
            if per == '':
                sheet.cell(2 + j, 2 + i).value = ''
            else:
                sheet.cell(2 + j, 2 + i).value = "%.1f" % (per)
    return table


def save_table():
    try:
        table.save(r'{}{}.xlsx'.format(project_path, table_result))
    except Exception as e:
        print(e)


def change_int_type(val0):
    val = ''
    for v in range(len(val0)):
        val += val0[v]
        if (len(val0) - v) % 3 == 1 and v != len(val0) - 1:
            val += ','
    return val


def get_ppt():
    tt = time.strftime("%m%d", time.localtime())
    myoffice.open_file(r'{}美素佳儿OUTPUT({}).pptx'.format(project_path,tt), template=r'{}{}'.format(project_path,ppt_template)).fill(r'{}{}.xlsx'.format(project_path,table_result)).save()
    myppt = Presentation(r'{}美素佳儿OUTPUT({}).pptx'.format(project_path,tt))
    for i in range(len(myppt.slides)):
        if i >1:
            for shape in myppt.slides[i].shapes:
                if shape.has_table:
                    for row in range(len(shape.table.rows)):
                        for col in range(len(shape.table.columns)):
                            if shape.table.cell(row,col).text_frame.text != '':
                                try:
                                    if len(re.findall(r'[-%]',shape.table.cell(row,col).text_frame.text)) == 2:
                                        shape.table.cell(row,col).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
                                    if len(re.findall(r'[+%]',shape.table.cell(row,col).text_frame.text)) == 2:
                                        shape.table.cell(row,col).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                                except:
                                    continue
    myppt.save(r'{}美素佳儿OUTPUT({}).pptx'.format(project_path,tt))

def set_env(report_date):
    global project_path,table_template,table_result,meisu,price_list,group_people,mouth,ppt_template,table,date1,date2,sql_list,newdata
    project_path = 'C:/Users/zeng.xiangyan/美素佳儿奶粉月报PPT数据更新/'
    table_template = 'source模板'
    table_result = 'source测试0816'
    ppt_template = '测试.pptx'
    meisu = ['行业|销售额(000元)','美素佳儿（金装+源悦+皇家）','皇家美素佳儿','旺玥','美素佳儿（金装）','美素佳儿（源悦）']
    mouth = [(datetime.strptime(report_date, '%Y-%m-%d')- relativedelta(months=4-i)).strftime('%Y%m') for i in range(1,4)]
    date1,date2,sql_list = config_date(report_date)
    get_source_model(report_date)
    newdata = importdata_new(sql_list)
    table=load_workbook(project_path + table_template + '.xlsx')

def run(start_date,end_date):
    global mat_list,table
    mat_list = ['MAT\nJun22','MAT\nJun23']
    set_env(end_date)
    process_data()
    get_summary()
    get_chart_data()
    get_table_data()
    total_price_chart()
    total_price_table()
    save_table()
    get_ppt()

if __name__ == "__main__":
    run('2023-11-01','2023-12-01')
